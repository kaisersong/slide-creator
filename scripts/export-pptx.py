#!/usr/bin/env python3
"""
export-pptx.py — Export slide-creator HTML presentations to editable PPTX.

No browser or Chromium required. Uses python-pptx + BeautifulSoup only.
Produces real PowerPoint files with editable text (not screenshots).

Usage:
    python export-pptx.py <presentation.html> [output.pptx]

Dependencies:
    pip install python-pptx beautifulsoup4
"""

import sys
import re
import base64
import tempfile
import os
from pathlib import Path


# ─── Dependency check ─────────────────────────────────────────────────────────

def check_deps():
    missing = []
    try:
        from bs4 import BeautifulSoup  # noqa
    except ImportError:
        missing.append("beautifulsoup4")
    try:
        from pptx import Presentation  # noqa
    except ImportError:
        missing.append("python-pptx")
    if missing:
        print(f"Missing dependencies. Install with:")
        print(f"  pip install {' '.join(missing)}")
        sys.exit(1)

check_deps()

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


# ─── Constants ────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)   # 16:9 widescreen
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.7)


# ─── Color helpers ────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str):
    """Parse #rgb or #rrggbb to RGBColor. Returns a dark grey on failure."""
    s = hex_str.strip().lstrip('#')
    if len(s) == 3:
        s = s[0]*2 + s[1]*2 + s[2]*2
    if len(s) == 6:
        try:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        except ValueError:
            pass
    return RGBColor(0x33, 0x33, 0x33)


def resolve_css_color(css_vars, key, fallback='#333333'):
    """
    Resolve a CSS variable name to RGBColor.
    Follows one level of var() reference (e.g. --card-bg: var(--bg-secondary)).
    """
    val = css_vars.get(key, fallback)
    # Follow one var() reference
    ref = re.match(r'var\(--([\w-]+)\)', val)
    if ref:
        val = css_vars.get(ref.group(1), fallback)
    # Extract the first hex color found
    m = re.search(r'#[0-9a-fA-F]{3,8}', val)
    return hex_to_rgb(m.group(0)) if m else hex_to_rgb(fallback)


# ─── CSS variable extraction ──────────────────────────────────────────────────

def extract_css_vars(html_content):
    """
    Extract all --variable: value; pairs from the first :root {} block.
    Returns a dict of variable_name -> raw_value_string.
    """
    vars = {}
    # Match :root { ... } (non-greedy, first match)
    m = re.search(r':root\s*\{([^}]+)\}', html_content, re.DOTALL)
    if m:
        for line in m.group(1).splitlines():
            pair = re.match(r'\s*--([\w-]+)\s*:\s*(.+?)\s*;', line)
            if pair:
                vars[pair.group(1)] = pair.group(2).strip()
    return vars


# ─── Font mapping ─────────────────────────────────────────────────────────────

# Maps web font names (lowercase substrings) to system fonts safe in PowerPoint
WEB_FONT_MAP = {
    'clash display':      'Arial',
    'archivo black':      'Arial Black',
    'archivo':            'Arial',
    'manrope':            'Calibri',
    'satoshi':            'Calibri',
    'syne':               'Arial',
    'space grotesk':      'Calibri',
    'plus jakarta sans':  'Calibri',
    'outfit':             'Calibri',
    'dm sans':            'Calibri',
    'work sans':          'Calibri',
    'nunito':             'Calibri',
    'ibm plex sans':      'Calibri',
    'cormorant garamond': 'Georgia',
    'cormorant':          'Georgia',
    'bodoni moda':        'Georgia',
    'fraunces':           'Georgia',
    'source serif':       'Georgia',
    'space mono':         'Courier New',
    'jetbrains mono':     'Courier New',
}

def map_font(raw_font_value):
    """
    Given a CSS font-family value like "'Clash Display', sans-serif",
    return the best matching system font name.
    """
    # Pull out the first font name (strip quotes)
    first = re.search(r"['\"]?([A-Za-z][^,'\"]+)['\"]?", raw_font_value or '')
    if not first:
        return 'Calibri'
    key = first.group(1).strip().lower()
    for pattern, mapped in WEB_FONT_MAP.items():
        if pattern in key:
            return mapped
    return 'Calibri'

def theme_fonts(css_vars):
    """Return (display_font, body_font) as system font names."""
    return (
        map_font(css_vars.get('font-display', '')),
        map_font(css_vars.get('font-body', '')),
    )


# ─── Shape helpers ────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color, no_line=True):
    """Add a filled rectangle with optional border removal."""
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.color.rgb = fill_color  # hide border by matching fill
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name, font_size, bold=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
    """Add a simple text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def add_bullets(slide, items, left, top, width, height,
                font_name, font_size, color=None):
    """Add a bulleted list. Each item gets a • prefix."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = font_name
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color
    return txb


# ─── Slide builder ────────────────────────────────────────────────────────────

def build_slide(prs, slide_elem, css_vars, html_dir):
    """
    Add one slide to the Presentation based on a .slide BeautifulSoup element.
    Detects layout from content and slide classes.
    """
    # — Theme colors
    bg      = resolve_css_color(css_vars, 'bg-primary',      '#1a1a1a')
    text    = resolve_css_color(css_vars, 'text-primary',     '#ffffff')
    dim     = resolve_css_color(css_vars, 'text-secondary',   '#aaaaaa')
    accent  = resolve_css_color(css_vars, 'accent',           '#00ffcc')
    card_bg = resolve_css_color(css_vars, 'bg-secondary',     '#111827')

    # — Theme fonts
    disp_font, body_font = theme_fonts(css_vars)

    # — Blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = bg

    # — Extract content
    classes   = slide_elem.get('class') or []
    h1        = slide_elem.find('h1')
    h2        = slide_elem.find('h2')
    heading   = h1 or h2
    title_txt = heading.get_text(strip=True) if heading else ''

    # Subtitle: first <p> or <h3> that isn't inside a list/code
    subtitle_txt = ''
    for tag in slide_elem.find_all(['p', 'h3'], recursive=True):
        if tag.find_parent(['ul', 'ol', 'pre', 'code']):
            continue
        t = tag.get_text(strip=True)
        if t and t != title_txt:
            subtitle_txt = t
            break

    # Bullets
    bullets = []
    ul = slide_elem.find(['ul', 'ol'])
    if ul:
        for li in ul.find_all('li', recursive=False):
            t = li.get_text(strip=True)
            if t:
                bullets.append(t)

    # Code
    pre  = slide_elem.find('pre')
    code = pre.get_text() if pre else (slide_elem.find('code') or None)
    if code and not isinstance(code, str):
        code = code.get_text()
    code_txt = (code or '').strip()

    # Images
    imgs = slide_elem.find_all('img')

    # — Layout detection
    is_title = 'title-slide' in classes or (h1 and not bullets and not code_txt)

    content_top = Inches(1.9)
    full_w = SLIDE_W - 2 * MARGIN

    def accent_bar(top=Inches(7.1), height=Pt(4)):
        add_rect(slide, Inches(0), top, SLIDE_W, height, accent)

    def title_block(font_size=32, top=Inches(0.5)):
        if title_txt:
            add_textbox(slide, title_txt, MARGIN, top, full_w, Inches(1.2),
                        disp_font, font_size, bold=True, color=text)
            # Small accent underline
            add_rect(slide, MARGIN, top + Inches(1.25), Inches(1.2), Pt(3), accent)

    # ── Title slide ───────────────────────────────────────────────────────────
    if is_title:
        if title_txt:
            add_textbox(slide, title_txt,
                        MARGIN, Inches(2.3), full_w, Inches(2.2),
                        disp_font, 44, bold=True, color=text, align=PP_ALIGN.CENTER)
        if subtitle_txt:
            add_textbox(slide, subtitle_txt,
                        MARGIN, Inches(4.7), full_w, Inches(1.0),
                        body_font, 20, color=dim, align=PP_ALIGN.CENTER)
        accent_bar()

    # ── Bullets ───────────────────────────────────────────────────────────────
    elif bullets:
        title_block()
        add_bullets(slide, bullets, MARGIN, content_top, full_w, Inches(5.0),
                    body_font, 19, color=text)

    # ── Code slide ────────────────────────────────────────────────────────────
    elif code_txt:
        title_block()
        code_box = slide.shapes.add_textbox(MARGIN, content_top, full_w, Inches(4.5))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = card_bg
        tf = code_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = code_txt
        run.font.name = 'Courier New'
        run.font.size = Pt(13)
        run.font.color.rgb = accent

    # ── Image slide ───────────────────────────────────────────────────────────
    elif imgs:
        title_block()
        img_added = False
        for img_elem in imgs[:1]:
            src = img_elem.get('src', '')
            try:
                if src.startswith('data:'):
                    # base64 embedded image — decode to temp file
                    header, b64data = src.split(',', 1)
                    ext = re.search(r'data:image/(\w+)', header)
                    suffix = f".{ext.group(1)}" if ext else '.png'
                    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
                    tmp.write(base64.b64decode(b64data))
                    tmp.close()
                    img_path = tmp.name
                    cleanup = True
                else:
                    img_path = str(html_dir / src)
                    cleanup = False

                if os.path.exists(img_path):
                    # Fit image in content area (max 10" wide, 4.5" tall)
                    slide.shapes.add_picture(
                        img_path,
                        MARGIN, content_top,
                        min(full_w, Inches(10)),
                        Inches(4.5)
                    )
                    img_added = True
                    if cleanup:
                        os.unlink(img_path)
            except Exception as e:
                print(f"    ⚠ Skipped image ({src[:40]}...): {e}")

        if not img_added and subtitle_txt:
            add_textbox(slide, subtitle_txt, MARGIN, content_top, full_w, Inches(4),
                        body_font, 20, color=dim)

    # ── Generic (title + text) ────────────────────────────────────────────────
    else:
        title_block(font_size=36, top=Inches(0.6))
        body_parts = []
        for tag in slide_elem.find_all(['p', 'h3', 'h4'], recursive=True):
            if tag.find_parent(['ul', 'ol', 'pre']):
                continue
            t = tag.get_text(strip=True)
            if t and t != title_txt and t != subtitle_txt:
                body_parts.append(t)
        body_txt = '\n\n'.join(body_parts[:4]) or subtitle_txt
        if body_txt:
            add_textbox(slide, body_txt, MARGIN, content_top, full_w, Inches(4.5),
                        body_font, 20, color=text)

    return slide


# ─── Main ─────────────────────────────────────────────────────────────────────

def export(html_path, output_path=None):
    html_path = Path(html_path).resolve()
    if not html_path.exists():
        print(f"Error: file not found: {html_path}")
        sys.exit(1)

    html_dir     = html_path.parent
    html_content = html_path.read_text(encoding='utf-8')
    soup         = BeautifulSoup(html_content, 'html.parser')
    css_vars     = extract_css_vars(html_content)

    slides_elems = soup.select('.slide')
    if not slides_elems:
        print("No .slide elements found in HTML.")
        sys.exit(1)

    total = len(slides_elems)
    print(f"Exporting {total} slides from: {html_path.name}")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, elem in enumerate(slides_elems, 1):
        title = (elem.find('h1') or elem.find('h2') or '').get_text(strip=True) if \
                (elem.find('h1') or elem.find('h2')) else f'Slide {i}'
        print(f"  [{i}/{total}] {title[:60]}")
        build_slide(prs, elem, css_vars, html_dir)

    out = output_path or str(html_path.with_suffix('.pptx'))
    prs.save(out)
    print(f"\n✓ Saved: {out}  ({total} slides)")
    return out


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(0)
    export(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
